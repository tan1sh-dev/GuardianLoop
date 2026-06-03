from pptx import Presentation

replacements = {
    'GuardianLoop is an autonomous, multi-agent': 'GuardianLoop is a paradigm shift in application security: an autonomous, multi-agent pipeline that manages the entire vulnerability lifecycle—from detection to verified remediation—with zero human intervention. Unlike conventional Static Application Security Testing (SAST) tools that merely flag issues, GuardianLoop actively patches the code and cryptographically or empirically proves the fix works before surfacing it, effectively bridging the gap between discovery and resolution.',
    'Static analysis tools like Semgrep and Bandit': 'Modern software development suffers from a broken, labor-intensive security workflow:\n• False Positive Fatigue: Traditional SAST tools dump raw, uncontextualized findings on engineers.\n• The "Unverified AI" Trap: AI coding assistants suggest fixes probabilistically, offering no guarantee that the patch neutralizes the exploit.\n• Persisting Vulnerabilities: Known vulnerabilities often persist in production despite heavy investments in tooling.',
    'Automate end-to-end vulnerability management': '• End-to-End Automation: Seamlessly orchestrate the pipeline from vulnerability discovery to a verified, deployed fix.\n• Context-Aware Triage: Enrich raw SAST findings with real-world threat intelligence via the NVD REST API.\n• Semantic Remediation: Generate highly accurate, logically reasoned patches using Chain-of-Thought (CoT) prompting.\n• Deterministic Validation: Guarantee patch effectiveness by executing the original exploit against the patched code inside a heavily restricted sandbox.\n• Auditable Reporting: Produce transparent, human-readable CI/CD artifacts.',
    '[Flowchart to elaborate': 'The system utilizes a stateful LangGraph architecture to route data through five specialized agents:\n1. Scout: Scans the codebase for vulnerabilities.\n2. Classifier: Prioritizes findings using global threat intel.\n3. Fixer: Drafts a semantic code patch.\n4. Red-Team (The Core Feedback Loop): Deploys a locked-down Docker sandbox to run the exploit. If the exploit succeeds, execution logs are fed back to the Fixer (up to 3 retries) until neutralized.\n5. Reporter: Compiles the successful verification into a structured audit artifact.',
    'The five agents run sequentially': 'The multi-agent orchestration successfully decouples complex reasoning tasks. The standout innovation is the Red-Team Agent. By forcing the LLM-generated patch to face a live exploit in a deterministic Docker environment, GuardianLoop converts probabilistic AI outputs into a binary pass/fail outcome.\nDuring testing, the iterative feedback loop drastically increased the success rate of complex patches. Only fixes that actively block the exploit are approved—meaning zero unverified code is ever shipped.',
    'GuardianLoop demonstrates that a complete': 'GuardianLoop proves that a complete, self-correcting vulnerability lifecycle is achievable today. By combining LLM-driven Chain-of-Thought reasoning with deterministic containerized sandboxing, we have eliminated the manual overhead of triage and patch validation. The system effectively upgrades security from a passive warning system to an autonomous self-healing ecosystem.',
    'Five-agent LangGraph pipeline': '• A robust, five-agent LangGraph pipeline bound by strict state-passing contracts.\n• Hardened, ephemeral Docker sandbox environments configured for Python and C++ exploit testing.\n• Structured, CI-ready output artifacts (run_summary.json, patches.json, report.md).\n• A fully mocked, zero-dependency test suite allowing rapid offline development and demonstration.'
}

try:
    prs = Presentation('e:\\GuardianLoop\\EL Poster draft 1.pptx')
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text
            for key, new_text in replacements.items():
                if key in text:
                    shape.text = new_text
    prs.save('e:\\GuardianLoop\\EL Poster draft 1_optimized.pptx')
    print('Successfully optimized pptx!')
except Exception as e:
    print('Error:', e)
